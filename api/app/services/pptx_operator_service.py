import os
import shutil

from pptx import Presentation


class PPTXOperatorService:
    def __init__(self, pptx_file_name, save_path):
        self.pptx_file_name = pptx_file_name
        self.pptx_file_path = os.path.join(
            os.path.dirname(__file__), "..", "static", "template", pptx_file_name
        )
        self.new_pptx_file_path = os.path.join(save_path, "output.pptx")
        self.save_path = save_path

    def insert_image(self, ai_image_file_name):
        # 元ファイル(pptx)をコピー
        shutil.copy(self.pptx_file_path, self.new_pptx_file_path)

        # コピーしたpptxファイルを読み込む
        presentation = Presentation(self.new_pptx_file_path)

        # 最初のスライド取得
        slide = presentation.slides[0]

        # スライドサイズを取得
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # 画像をスライドに追加（最背面にするには最初に追加）
        ai_image_file_path = os.path.join(self.save_path, ai_image_file_name)
        slide.shapes.add_picture(
            ai_image_file_path, 0, 0, width=slide_width, height=slide_height
        )

        # Indexの順序を更新
        spTree = slide.shapes._spTree
        background_shape = slide.shapes[len(slide.shapes) - 1]

        # 他のshapeを一旦削除
        other_shapes = [s for s in slide.shapes if s != background_shape]
        for shape in other_shapes:
            spTree.remove(shape._element)

        # 他のshapeを再追加（順番最後になるので前面になる）
        for shape in other_shapes:
            spTree.append(shape._element)

        # PPTX保存
        presentation.save(self.new_pptx_file_path)

        return

    def export_image(self):
        pass
